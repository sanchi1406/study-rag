from pptx import Presentation


def load_ppt(ppt_path: str) -> str:
    """
    Extract text from a PowerPoint presentation.
    """

    prs = Presentation(ppt_path)

    text = []

    for slide in prs.slides:

        for shape in slide.shapes:

            if hasattr(shape, "text"):

                if shape.text.strip():

                    text.append(shape.text)

    return "\n".join(text)